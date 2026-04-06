import json
import os
import re
import subprocess
import tempfile
import urllib.request
import urllib.error
from http.server import BaseHTTPRequestHandler
from datetime import datetime

NOTION_API_KEY = os.environ.get("NOTION_API_KEY", "")
NOTION_DATABASE_ID = os.environ.get("NOTION_DATABASE_ID", "036bdd7a61694c0e95450a26984e84c4")
GOOGLE_CLIENT_ID = os.environ.get("GOOGLE_CLIENT_ID", "")
GOOGLE_CLIENT_SECRET = os.environ.get("GOOGLE_CLIENT_SECRET", "")
GOOGLE_REFRESH_TOKEN = os.environ.get("GOOGLE_REFRESH_TOKEN", "")

def notion_request(method, path, payload=None):
    url = f"https://api.notion.com/v1{path}"
    headers = {
        "Authorization": f"Bearer {NOTION_API_KEY}",
        "Content-Type": "application/json",
        "Notion-Version": "2022-06-28"
    }
    data = json.dumps(payload).encode() if payload else None
    req = urllib.request.Request(url, data=data, headers=headers, method=method)
    with urllib.request.urlopen(req) as resp:
        return json.loads(resp.read())

def get_google_access_token():
    """Get a fresh Google access token using the refresh token."""
    url = "https://oauth2.googleapis.com/token"
    payload = {
        "client_id": GOOGLE_CLIENT_ID,
        "client_secret": GOOGLE_CLIENT_SECRET,
        "refresh_token": GOOGLE_REFRESH_TOKEN,
        "grant_type": "refresh_token"
    }
    req = urllib.request.Request(
        url,
        data=json.dumps(payload).encode(),
        headers={"Content-Type": "application/json"},
        method="POST"
    )
    with urllib.request.urlopen(req) as resp:
        return json.loads(resp.read())["access_token"]

def google_request(method, url, token, payload=None, content_type="application/json"):
    headers = {
        "Authorization": f"Bearer {token}",
        "Content-Type": content_type
    }
    data = json.dumps(payload).encode() if payload and content_type == "application/json" else payload
    req = urllib.request.Request(url, data=data, headers=headers, method=method)
    with urllib.request.urlopen(req) as resp:
        return json.loads(resp.read())

def find_notion_candidate(candidate_name):
    try:
        data = notion_request("POST", f"/databases/{NOTION_DATABASE_ID}/query", {
            "filter": {"property": "title", "title": {"contains": candidate_name}}
        })
        results = data.get("results", [])
        target = candidate_name.lower().strip()

        for page in results:
            name_prop = (page.get("properties", {}).get("\ufeffName") or
                        page.get("properties", {}).get("Name", {}))
            page_name = "".join(
                t.get("plain_text", "") for t in name_prop.get("title", [])
            ).lower().strip()
            if page_name == target:
                return page["id"], page

        best, best_score = None, 0
        for page in results:
            name_prop = (page.get("properties", {}).get("\ufeffName") or
                        page.get("properties", {}).get("Name", {}))
            page_name = "".join(
                t.get("plain_text", "") for t in name_prop.get("title", [])
            ).lower().strip()
            score = sum(1 for w in target.split() if w in page_name.split())
            if score > best_score:
                best_score, best = score, page

        if best:
            return best["id"], best
    except Exception:
        pass
    return None, None

def get_candidate_data(page):
    """Extract all relevant candidate data from a Notion page."""
    props = page.get("properties", {})

    def get_rich_text(prop_name):
        prop = props.get(prop_name, {})
        return "".join(t.get("plain_text", "") for t in prop.get("rich_text", []))

    def get_url(prop_name):
        return props.get(prop_name, {}).get("url", "") or ""

    def get_select(prop_name):
        s = props.get(prop_name, {}).get("select")
        return s.get("name", "") if s else ""

    def get_status(prop_name):
        s = props.get(prop_name, {}).get("status")
        return s.get("name", "") if s else ""

    def get_multi_select(prop_name):
        items = props.get(prop_name, {}).get("multi_select", [])
        return ", ".join(i.get("name", "") for i in items)

    name_prop = (props.get("\ufeffName") or props.get("Name", {}))
    full_name = "".join(t.get("plain_text", "") for t in name_prop.get("title", [])).strip()

    notes = get_rich_text("Notes")
    linkedin = get_url("LinkedIn")
    location = get_multi_select("Current Location")
    current_company = get_rich_text("Current Company")
    stage = get_status("Stage")

    # Parse notes sections
    def extract_section(notes, header):
        pattern = rf'{header}\s*\n(.*?)(?=\n[A-Z][A-Z ]+\n|$)'
        match = re.search(pattern, notes, re.DOTALL | re.IGNORECASE)
        return match.group(1).strip() if match else ""

    basics = extract_section(notes, "BASICS")
    strong_points = extract_section(notes, "STRONG POINTS")
    potential_challenges = extract_section(notes, "POTENTIAL CHALLENGES")
    compensation = extract_section(notes, "COMPENSATION")

    return {
        "full_name": full_name,
        "location": location,
        "current_company": current_company,
        "linkedin": linkedin,
        "stage": stage,
        "basics": basics,
        "strong_points": strong_points,
        "potential_challenges": potential_challenges,
        "compensation": compensation,
        "notes": notes
    }

def find_drive_folder(token, parent_id, folder_name):
    """Find a folder by name inside a parent folder."""
    query = f"name = '{folder_name}' and mimeType = 'application/vnd.google-apps.folder' and '{parent_id}' in parents and trashed = false"
    url = f"https://www.googleapis.com/drive/v3/files?q={urllib.parse.quote(query)}&fields=files(id,name)"
    import urllib.parse
    result = google_request("GET", url, token)
    files = result.get("files", [])
    return files[0]["id"] if files else None

def find_presentation_in_folder(token, folder_id):
    """Find the Google Slides presentation in a folder."""
    import urllib.parse
    query = f"mimeType = 'application/vnd.google-apps.presentation' and '{folder_id}' in parents and trashed = false"
    url = f"https://www.googleapis.com/drive/v3/files?q={urllib.parse.quote(query)}&fields=files(id,name,webViewLink)"
    result = google_request("GET", url, token)
    files = result.get("files", [])
    return files[0] if files else None

def build_slide_report(token, presentation_id, candidate):
    """Download PPTX, duplicate template slide, populate, re-upload."""
    import copy
    import urllib.parse

    # Download as PPTX
    export_url = f"https://www.googleapis.com/drive/v3/files/{presentation_id}/export?mimeType=application/vnd.openxmlformats-officedocument.presentationml.presentation"
    req = urllib.request.Request(export_url, headers={"Authorization": f"Bearer {token}"})
    with urllib.request.urlopen(req) as resp:
        pptx_bytes = resp.read()

    # Save to temp file
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as f:
        f.write(pptx_bytes)
        tmp_in = f.name

    tmp_out = tmp_in.replace('.pptx', '_updated.pptx')

    # Run python-pptx script
    name_parts = candidate["full_name"].split()
    first_name = name_parts[0] if name_parts else candidate["full_name"]
    last_name = " ".join(name_parts[1:]) if len(name_parts) > 1 else ""

    script = f'''
import copy, sys
from pptx import Presentation
from pptx.util import Pt
from lxml import etree

prs = Presentation("{tmp_in}")

# Find template slide (last slide with placeholder text)
template_idx = len(prs.slides) - 1
for i, slide in enumerate(prs.slides):
    for shape in slide.shapes:
        if shape.has_text_frame:
            text = shape.text_frame.text.upper()
            if "CANDIDATE" in text or "NAME" in text or "BASICS" in text:
                template_idx = i
                break

template_slide = prs.slides[template_idx]

# Deep copy the template slide
template_sp_tree = template_slide.shapes._spTree
slide_layout = template_slide.slide_layout
new_slide = prs.slides.add_slide(slide_layout)
new_sp_tree = new_slide.shapes._spTree

# Remove auto-added shapes
for i in range(len(new_sp_tree) - 1, -1, -1):
    new_sp_tree.remove(new_sp_tree[i])

# Copy all shapes from template
for child in template_sp_tree:
    new_sp_tree.append(copy.deepcopy(child))

# Helper to set text preserving formatting
def set_text(para, text):
    if not para.runs:
        return
    para.runs[0].text = text
    for r in para.runs[1:]:
        r.text = ""

candidate_data = {{
    "first_name": {repr(first_name)},
    "last_name": {repr(last_name)},
    "location": {repr(candidate["location"])},
    "linkedin": {repr(candidate["linkedin"])},
    "stage": {repr(candidate["stage"])},
    "basics": {repr(candidate["basics"])},
    "strong_points": {repr(candidate["strong_points"])},
    "potential_challenges": {repr(candidate["potential_challenges"])},
    "compensation": {repr(candidate["compensation"])},
}}

# Populate shapes
for shape in new_slide.shapes:
    if not shape.has_text_frame:
        continue
    full_text = shape.text_frame.text.upper()
    paras = shape.text_frame.paragraphs

    if "CANDIDATE" in full_text and "NAME" in full_text:
        if len(paras) >= 2:
            set_text(paras[0], candidate_data["first_name"])
            set_text(paras[1], candidate_data["last_name"])
        elif paras:
            set_text(paras[0], candidate_data["first_name"] + " " + candidate_data["last_name"])

    elif "BASICS" in full_text:
        for i, para in enumerate(paras):
            pt = para.text.upper().strip()
            if "BASICS" in pt:
                if i + 1 < len(paras):
                    set_text(paras[i+1], candidate_data["basics"])
            elif "STRONG POINTS" in pt or "STRONG" in pt:
                if i + 1 < len(paras):
                    set_text(paras[i+1], candidate_data["strong_points"])
            elif "POTENTIAL" in pt or "CHALLENGES" in pt:
                if i + 1 < len(paras):
                    set_text(paras[i+1], candidate_data["potential_challenges"])
            elif "COMPENSATION" in pt:
                if i + 1 < len(paras):
                    set_text(paras[i+1], candidate_data["compensation"])
            elif "STATUS" in pt:
                set_text(para, "Status: " + candidate_data["stage"])

    elif "LOCATION" in full_text or "CITY" in full_text:
        set_text(paras[0], candidate_data["location"])

    elif "LINKEDIN" in full_text:
        set_text(paras[0], candidate_data["linkedin"] or "LinkedIn")

prs.save("{tmp_out}")
print("OK")
'''

    result = subprocess.run(['python3', '-c', script], capture_output=True, text=True, timeout=60)
    if result.returncode != 0 or "OK" not in result.stdout:
        raise Exception(f"Slide generation failed: {result.stderr[:300]}")

    # Re-upload to Drive
    with open(tmp_out, 'rb') as f:
        pptx_data = f.read()

    # Upload as new file then convert, or update existing
    metadata = json.dumps({
        "name": "Candidate Report",
        "mimeType": "application/vnd.google-apps.presentation"
    }).encode()

    boundary = "MultipartBoundary"
    body = (
        f"--{boundary}\r\nContent-Type: application/json\r\n\r\n".encode() +
        metadata +
        f"\r\n--{boundary}\r\nContent-Type: application/vnd.openxmlformats-officedocument.presentationml.presentation\r\n\r\n".encode() +
        pptx_data +
        f"\r\n--{boundary}--".encode()
    )

    # Actually, just update the existing presentation by uploading the PPTX back
    upload_url = f"https://www.googleapis.com/upload/drive/v3/files/{presentation_id}?uploadType=media"
    req = urllib.request.Request(
        upload_url,
        data=pptx_data,
        headers={
            "Authorization": f"Bearer {token}",
            "Content-Type": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            "Content-Length": str(len(pptx_data))
        },
        method="PATCH"
    )
    with urllib.request.urlopen(req) as resp:
        json.loads(resp.read())

    # Cleanup
    try:
        os.unlink(tmp_in)
        os.unlink(tmp_out)
    except:
        pass

    return True

import urllib.parse

class Handler(BaseHTTPRequestHandler):
    def do_POST(self):
        try:
            cl = int(self.headers.get('Content-Length', 0))
            body = json.loads(self.rfile.read(cl))

            candidate_name = body.get('candidate_name', '').strip()
            client = body.get('client', '').strip()
            role = body.get('role', '').strip()

            if not candidate_name:
                return self._err(400, 'Candidate name is required.')
            if not client or not role:
                return self._err(400, 'Client and role are required.')

            # 1. Find candidate in Notion
            page_id, page = find_notion_candidate(candidate_name)
            if not page_id:
                return self._err(404, f'"{candidate_name}" not found in Notion. Check the name matches exactly.')

            # 2. Get candidate data
            candidate = get_candidate_data(page)

            # 3. Get Google access token
            token = get_google_access_token()

            # 4. Find client folder in Drive
            drive_root_url = "https://www.googleapis.com/drive/v3/files?q=mimeType%3D'application%2Fvnd.google-apps.folder'+and+name%3D'" + urllib.parse.quote(client) + "'+and+trashed%3Dfalse&fields=files(id,name)"
            result = google_request("GET", drive_root_url, token)
            client_folders = result.get("files", [])
            if not client_folders:
                return self._err(404, f'Client folder "{client}" not found in Google Drive.')
            client_folder_id = client_folders[0]["id"]

            # 5. Find role folder inside client folder
            role_folder_id = find_drive_folder(token, client_folder_id, role)
            if not role_folder_id:
                return self._err(404, f'Role folder "{role}" not found inside "{client}" in Google Drive.')

            # 6. Find presentation in role folder
            presentation = find_presentation_in_folder(token, role_folder_id)
            if not presentation:
                return self._err(404, f'No Google Slides presentation found in "{client} → {role}".')

            # 7. Build and upload the slide
            build_slide_report(token, presentation["id"], candidate)

            self._ok({
                'success': True,
                'deck_url': presentation.get("webViewLink", f"https://docs.google.com/presentation/d/{presentation['id']}")
            })

        except urllib.error.HTTPError as e:
            self._err(500, f'API error: {e.read().decode()[:200]}')
        except Exception as e:
            self._err(500, str(e))

    def _ok(self, data):
        self._json(200, data)

    def _err(self, code, msg):
        self._json(code, {'error': msg})

    def _json(self, status, data):
        body = json.dumps(data).encode()
        self.send_response(status)
        self.send_header('Content-Type', 'application/json')
        self.send_header('Content-Length', str(len(body)))
        self.send_header('Access-Control-Allow-Origin', '*')
        self.end_headers()
        self.wfile.write(body)

    def log_message(self, *args): pass
